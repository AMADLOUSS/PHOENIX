import pandas as pd
import numpy as np
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from openpyxl import load_workbook
from openpyxl.utils.dataframe import dataframe_to_rows

def read_excel_mapping(mapping_file):
    xls = pd.ExcelFile(mapping_file)
    catalogue_df = pd.read_excel(xls, sheet_name=0)
    mapping_df = pd.read_excel(xls, sheet_name=1)
    return catalogue_df, mapping_df

def read_pdf(file_path):
    text = ""
    if os.path.exists(file_path):
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
    return text

def read_pptx(file_path):
    text = ""
    if os.path.exists(file_path):
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error reading {file_path}: {e}")
    return text

def read_xlsx(file_path):
    return pd.read_excel(file_path)

def merge_input_files(input_files):
    all_data = []

    for file_path in input_files:
        if file_path.endswith('.xlsx'):
            data_df = read_xlsx(file_path)
            all_data.append(data_df.values.flatten())
        elif file_path.endswith('.pdf'):
            text = read_pdf(file_path)
            all_data.append([text])
        elif file_path.endswith('.pptx'):
            text = read_pptx(file_path)
            all_data.append([text])
        else:
            continue

    # Flatten and combine all data into one column
    combined_data = [item.strip() for sublist in all_data for item in sublist if isinstance(item, str)]
    combined_data_df = pd.DataFrame({'Combined': combined_data})
    
    return combined_data_df

def process_files(input_files, mapping_file):
    catalogue_df, mapping_df = read_excel_mapping(mapping_file)

    # Merge input files into one DataFrame
    data_df = merge_input_files(input_files)
    
    matched_columns = {col: False for col in catalogue_df.columns}

    # Get the second line (index 0) of the catalogue DataFrame for comparison
    second_line_values = catalogue_df.iloc[0]

    # Clean and split combined values into lists
    data_df['Combined'] = data_df['Combined'].apply(lambda x: [item.strip() for item in x.split('\n')] if isinstance(x, str) else x)

    # Explode the lists into separate rows
    data_df = data_df.explode('Combined').reset_index(drop=True)
    
    for item in data_df['Combined']:
        # Check if the item is in the "Original" column of the mapping DataFrame
        if item in mapping_df['Original'].values:
            # Get the corresponding "Transcoded" value
            transcoded_value = mapping_df.loc[mapping_df['Original'] == item, 'Transcoded'].values[0]
            
            for col in catalogue_df.columns:
                if transcoded_value == second_line_values[col] and col == item:
                    print(transcoded_value, col)
                    catalogue_df[col] = catalogue_df[col].apply(lambda x: 1 if pd.isna(x) else x)

    return catalogue_df

def save_to_excel(catalogue_df, mapping_file, output_file):
    if catalogue_df.empty:
        print("Warning: The catalogue DataFrame is empty! Not saving to output file.")
        return

    # Load original workbook to retain formatting
    original_wb = load_workbook(mapping_file)
    original_sheet = original_wb.worksheets[0]

    # Clear the original sheet contents while maintaining the formatting
    for row in original_sheet.iter_rows(min_row=1, max_col=original_sheet.max_column, max_row=original_sheet.max_row):
        for cell in row:
            cell.value = None

    # Write updated DataFrame to the original sheet
    for r_idx, r in enumerate(dataframe_to_rows(catalogue_df, index=False, header=True)):
        for c_idx, value in enumerate(r):
            original_sheet.cell(row=r_idx + 1, column=c_idx + 1, value=value)

    original_wb.save(output_file)
    print(f"Workbook saved as '{output_file}'")

# Example usage
if __name__ == "__main__":
    input_files = ['data1.xlsx', 'data2.pdf', 'data3.pptx']  # Update these to your input files
    mapping_file = 'mapping_file.xlsx'  # Your mapping file
    output_file = 'updated_catalogue.xlsx'  # New output file

    result_df = process_files(input_files, mapping_file)
    save_to_excel(result_df, mapping_file, output_file)

    if not result_df.empty:
        print("Results saved to:", output_file)
    else:
        print("No results to save.")
