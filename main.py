import pandas as pd
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from openpyxl import load_workbook
from openpyxl.utils.dataframe import dataframe_to_rows
from transformers import AutoTokenizer, AutoModel
import torch
import warnings

warnings.simplefilter("ignore", FutureWarning)

# Load the pretrained model and tokenizer
model_name = 'sentence-transformers/all-MiniLM-L6-v2'
tokenizer = AutoTokenizer.from_pretrained(model_name)
model = AutoModel.from_pretrained(model_name)

def read_file(file_path):
    """Read a file based on its extension."""
    ext = os.path.splitext(file_path)[1]
    if ext == '.xlsx':
        return pd.read_excel(file_path)
    elif ext == '.pdf':
        return read_pdf(file_path)
    elif ext == '.pptx':
        return read_pptx(file_path)
    else:
        print(f"Unsupported file type: {ext}")
        return None

def read_pdf(file_path):
    """Extract text from a PDF file."""
    if not os.path.exists(file_path):
        return ""
    text = ""
    with open(file_path, 'rb') as file:
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def read_pptx(file_path):
    """Extract text from a PPTX file."""
    if not os.path.exists(file_path):
        return ""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return text

def merge_input_files(input_files):
    """Merge text from various input files into a single DataFrame."""
    all_data = []
    for file_path in input_files:
        data = read_file(file_path)
        if isinstance(data, pd.DataFrame):
            all_data.extend(data.values.flatten())
        elif isinstance(data, str):
            all_data.append(data)
    
    return pd.DataFrame({'Combined': [item.strip() for item in all_data if isinstance(item, str)]})

def get_similarity(text1, text2):
    """Calculate cosine similarity between two texts."""
    inputs = tokenizer([text1, text2], padding=True, truncation=True, return_tensors='pt')
    with torch.no_grad():
        embeddings = model(**inputs).last_hidden_state.mean(dim=1)
    return torch.nn.functional.cosine_similarity(embeddings[0], embeddings[1], dim=0).item()

def is_contextually_similar(word1, word2, threshold=0.5):
    """Check if two words are contextually similar based on a threshold."""
    if pd.isna(word1) or pd.isna(word2):
        return False
    return get_similarity(str(word1), str(word2)) > threshold

def process_files(input_files, mapping_file):
    """Process input files and update the catalogue DataFrame."""
    catalogue_df, mapping_df = pd.read_excel(mapping_file, sheet_name=None).values()
    data_df = merge_input_files(input_files)

    # Explode combined values into separate rows
    data_df['Combined'] = data_df['Combined'].str.split('\n')
    data_df = data_df.explode('Combined').reset_index(drop=True)

    for item in data_df['Combined']:
        for original in mapping_df['Original']:
            if is_contextually_similar(item, original):
                item = original  
                transcoded_value = mapping_df.loc[mapping_df['Original'] == original, 'Transcoded'].values[0]
                for col in catalogue_df.columns:
                    if transcoded_value == catalogue_df.iloc[0][col] and col == item:
                        catalogue_df[col] = catalogue_df[col].apply(lambda x: 1 if pd.isna(x) else x)
                        print(f"Matched '{original}' with '{col}'")

    return catalogue_df

def save_to_excel(catalogue_df, mapping_file, output_file):
    """Save the updated catalogue DataFrame to an Excel file."""
    if catalogue_df.empty:
        print("Warning: The catalogue DataFrame is empty! Not saving to output file.")
        return

    original_wb = load_workbook(mapping_file)
    original_sheet = original_wb.worksheets[0]

    # Clear original sheet contents while retaining formatting
    for row in original_sheet.iter_rows():
        for cell in row:
            cell.value = None

    # Write updated DataFrame to the original sheet
    for r_idx, row in enumerate(dataframe_to_rows(catalogue_df, index=False, header=True)):
        for c_idx, value in enumerate(row):
            original_sheet.cell(row=r_idx + 1, column=c_idx + 1, value=value)

    original_wb.save(output_file)
    print(f"Workbook saved as '{output_file}'")

# Example usage
if __name__ == "__main__":
    input_files = ['data1.xlsx', 'data2.pdf', 'data3.pptx'] 
    mapping_file = 'mapping_file.xlsx'  
    output_file = 'updated_catalogue.xlsx' 

    result_df = process_files(input_files, mapping_file)
    save_to_excel(result_df, mapping_file, output_file)

    if not result_df.empty:
        print("Results saved to:", output_file)
    else:
        print("No results to save.")
